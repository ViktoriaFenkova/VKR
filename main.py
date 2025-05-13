from click import password_option
from streamlit import button

from Functions import выбор_шаблона_streamlit, заполнение_шаблона_streamlit
import json
import streamlit as st


import re
from docx import Document


from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO

from PIL import Image


# Конфигурация страницы
st.set_page_config(
    page_title="ПВК-Конструктор",
    layout="wide",
    initial_sidebar_state="expanded",
    page_icon=Image.open("logo.png")
)

logo = Image.open("logo.png")# добавляю файл с логотипом

st.sidebar.image(logo, width=120)# добавляю логотип для отображения на каждой странице приложения внизу
st.sidebar.markdown("### ПВК-Конструктор")

# Добавляю логотип для отображения сверху, для этого создаю две колонки
col1, col2 = st.columns([1, 3])  # Изображение в первой колонке, текст — во второй
with col1:
    st.image(logo, width=400)
with col2:
    st.markdown("## ПВК-Конструктор")
    st.markdown(
        '<p style="color: gray;">централизованное решение для эффективного управления рисками в ПОД/ФТ в вашей организации</p>',
        unsafe_allow_html=True
    )

# Страница по умолчанию
if "page" not in st.session_state: #st.session_state - хранилище всех переменных, которые мы туда положили и они должны сохраняться при взаимодействии со страницей
    st.session_state.page = "page_main" # назначаю Главную страницу по умолчанию при входе на платформу

# Содержимое главной страницы
page_main = st.sidebar.button("📎 Главная")
if page_main == True:
    st.session_state.page = "page_main"

# Содержимое страницы с разделом Конструктор ПВК
user_regime = st.sidebar.button("🔧 Конструктор ПВК")
print(user_regime)
if user_regime == True:
    st.session_state.page = "page_PVK" # при каждом нажатии на кнопку переменная перезаписывается
    st.title("Конструктор ПВК")
print(st.session_state.page)

# Содержимое страницы с разделом Редактор Шаблонов
if "users_db" not in st.session_state:
    # Два режима: Администратор (admin) и пользователь (user)
    st.session_state.users_db = {
        "admin": {"admin": True, "password": "admin"},
        "user": {"admin": False, "password": "user"}
    }

    st.session_state.username = None
    st.session_state.password = None

if st.session_state.username in st.session_state.users_db and (st.session_state.users_db) [st.session_state.username] ["admin"]: # == True
    admin_page = st.sidebar.button("📌 Страничка админа")
    print(admin_page)
    if admin_page == True:
        st.session_state.page = "admin_page"

# Содержимое страницы с разделом Обучение по ПОД/ФТ
presentation_page = st.sidebar.button("📚 Обучение по ПОД/ФТ")
if presentation_page == True:
    st.session_state.page = "presentation_page"

# Содержимое страницы с разделом Новости ПОД/ФТ
news_page = st.sidebar.button("📢 Новости ПОД/ФТ")
print(news_page)
if news_page == True:
    st.session_state.page = "news_page"

шаблон_ПВК = """
templates = { 
    'ПВК': { - словарь (с 10 по 13 стр) ключ - название шаблона
        'template': "./Data/Templates/PVK Template.docx", -  - параметр словаря "ПВК" , путь к файлу шаблона
        'parameters': {
            'client': 'наименование компании'
            },
        'description': 'Описание ПВК для Депозитария'
    },
    'ПВК для ИП': {
        'template': './Data/Templates/ПВК для ИП.docx',
        'parameters': {
            'client': 'наименование ИП'
            },
        'description': 'Описание ПВК для ИП'
    }}
"""


with open("./Data/формы_шаблонов.json", "r") as templates_file: #with использую, чтобы открыть файл, прочитать его и закрыть
    templates = json.load(templates_file)

# Режим пользователя по умолчанию


# Функция для регистрации пользователя
def register_user(username, password):
    if username in st.session_state.users_db:
        st.error("Пользователь с таким именем уже существует.")
    else:
        st.session_state.users_db[username] = {"admin" : False, # любой пользователь - админ
                                               "password": password}
        st.success("Регистрация прошла успешно!")

# Функция для входа пользователя
def login_user(username, password):
    if username in st.session_state.users_db and st.session_state.users_db[username]["password"] == password:
        st.success("Вход выполнен успешно!")
    else:
        st.error("Неверное имя пользователя или пароль.")

if (not (st.session_state.username in st.session_state.users_db and st.session_state.users_db[st.session_state.username] ["password"] == st.session_state.password) and
    st.session_state.page != "page_main" and st.session_state.page != "news_page"):
    st.title("Форма регистрации и входа") # Заголовок формы регистрации и входа

    option = st.selectbox("Выберите действие", ["Регистрация", "Вход"])

    if option == "Регистрация":
        st.header("Регистрация")
        st.session_state.username = st.text_input("Имя пользователя")
        st.session_state.password = st.text_input("Пароль", type="password")

        if st.button("Зарегистрироваться"):
            register_user( st.session_state.username, st.session_state.password)

    elif option == "Вход":
        st.header("Вход")
        st.session_state.username = st.text_input("Имя пользователя")
        st.session_state.password = st.text_input("Пароль", type="password")

        if st.button("Войти"):
            login_user( st.session_state.username, st.session_state.password)

print(st.session_state.username)
print(st.session_state.users_db)


# Функция для извлечения параметров из шаблона
def extract_parameters_from_docx(docx_file):
    # Загружаем документ
    doc = Document(docx_file)

    # Создаю список для хранения параметров
    parameters = set()

    # Прохожу по всем параграфам в документе
    for para in doc.paragraphs:
        # Использую регулярное выражение для поиска параметров в фигурных скобках
        matches = re.findall(r'\{(.*?)\}', para.text)
        parameters.update(matches)  # Добавляю найденные параметры в множество

    return list(parameters), doc # возврат параметров

with open("./Data/news.json", "r") as news_file:
    news = json.load(news_file)

#st.session_state - хранение переменных в рамках одной сессии
# Наполнение главной страницы
if st.session_state.page == "page_main":
    st.subheader("ПВК-Конструктор — всё для работы по ПОД/ФТ в одном месте")
    st.markdown(
        '<p style="color: gray;">Добро пожаловать в профессиональный инструмент для специалистов в сфере внутреннего контроля, комплаенса и ПОД/ФТ. Платформа объединяет в себе ключевые функции, необходимые для соблюдения требований законодательства и повышения эффективности работы</p>',
        unsafe_allow_html=True
    )

    st.markdown("""
    
        ### 
        

        ---

        #### 🔧 **Конструктор Правил Внутреннего Контроля**
        Создавайте, редактируйте и актуализируйте **Правила внутреннего контроля** вашей компании в соответствии с актуальными нормативными требованиями.  
        С Конструктором вы легко адаптируете правила под специфику вашей деятельности.

        ---

        #### 📚 **Обучение сотрудников по ПОД/ФТ**
        Генерируйте **обучающие презентации** с текстами и изображениями для повышения осведомлённости персонала по вопросам ПОД/ФТ и исполнения требований законодательства по обучению сотрудников.  
        Конструктор презентаций позволяет просто и быстро сформировать наглядные обучающие материалы.

        ---

        #### 📰 **Новости ПОД/ФТ**
        Следите за последними изменениями в сфере ПОД/ФТ для своевременного внедрения обновлений в процессы и обучения сотрудников организации.

        ---
        """)


# Наполнение страницы с новостями
elif st.session_state.page == "news_page":
    st.write("Актуальные новости в сфере ПОД/ФТ и финансового мониторинга")
    st.title("Новости")

    # Пример статичных новостей
    новости ="""{
            "title": "Новость 1: Обновление функционала",
            "content": "Мы выпустили обновление, которое добавляет новые возможности для пользователей. Теперь вы можете генерировать отчёты в новом формате.",
            "date": "2025-04-25",
        },
        {
            "title": "Новость 2: Исправления ошибок",
            "content": "В следующем обновлении мы исправим несколько критических багов, которые могли мешать работе приложения.",
            "date": "2025-04-20",
        },
        {
            "title": "Новость 3: Запуск нового сервиса",
            "content": "Мы рады анонсировать запуск нового сервиса для автоматизации работы с данными. Узнайте все подробности в нашем блоге.",
            "date": "2025-04-15",
        },
        ]"""

    # Сортировка новостей по дате
    news_list = []
    for заголовок, текст_новости in news.items():
        news_list.append({
            'title': заголовок,
            'content': текст_новости['content'],
            'date': текст_новости['date']
        })
    news_list.sort(key=lambda x: x["date"], reverse=True)

    # Отображение новостей
    for news_item in news_list:
        st.subheader(news_item["title"])
        st.write(f"**Дата:** {news_item['date']}")
        st.write(news_item["content"])
        st.markdown("---")

# Работа с шаблонами
if st.session_state.username in st.session_state.users_db and st.session_state.users_db[st.session_state.username]["password"] == st.session_state.password:
    if st.session_state.page== "page_PVK":
        template_name = выбор_шаблона_streamlit(templates) # вызов функции с параметрами tempiates, и после этого функция возвращает результат и он записывается в переменную template_name

        template_dict = templates[template_name] # template -это словарь, ключи в этом словаре - наименования шаблонов, template_name-наименование конкретного шаблона, выбранного пользователем
        template_path = template_dict['template'] # в template_path записывается 'template': './Data/Templates/ПВК для ИП.docx',
    # with open(template_path, 'r', encoding='utf-8') as template_file:
    # template = template_file.read()

        заполненный_шаблон = заполнение_шаблона_streamlit(template_dict, template_path)
        print(заполненный_шаблон) # самопроверка (но можно сделать через чекпоинт - черз дебак режим)
        with open("шаблон_клиента.docx", "w") as шаблон_файл:
            шаблон_файл.write(заполненный_шаблон)# retern - только с не готовыми функциями
        with open("шаблон_клиента.docx", "r") as шаблон_файл:
            st.download_button("Скачать документ", шаблон_файл, file_name= template_name +".docx")
            # Страница администратора
    elif st.session_state.page == "admin_page":
        st.title("Страница админа")
        действие = st.selectbox("Действия", ["Добавить шаблон", "Удалить шаблоны", "Добавить новость", "Удалить новость"])
        if действие == "Добавить шаблон":
            template_name = st.text_input("Название шаблона")
            template_description = st.text_area("Описание шаблона")

            # Загрузка шаблона
            uploaded_file = st.file_uploader("Загрузите .docx шаблон", type=["docx"])

            # Словарь для хранения описаний параметров
            parameter_descriptions = {}
            parameters = []

            if uploaded_file:
                # Извлекаю параметры из загруженного шаблона
                parameters, doc = extract_parameters_from_docx(uploaded_file)

                if parameters:
                    st.markdown("### Параметры шаблона")


                    # Для каждого параметра запрашиваю описание
                    for param in parameters:
                        description = st.text_input(f"Описание для параметра: {param}", key=f"param_desc_{param}")
                        if description:
                            parameter_descriptions[param] = description
                else:
                    st.warning("В шаблоне не найдены параметры.")
            # Кнопка сохранения шаблона
            сохранение_шаблона = st.button("Сохранить шаблон")
            if сохранение_шаблона == True:
                if not template_name:
                    st.warning("Название шаблона не заполнено")
                    st.error("Шаблон не сохранен")
                elif not template_description:
                    st.warning("Описание не заполнено")
                    st.error("Шаблон не сохранен")
                elif not uploaded_file:
                    st.warning("Файл не загружен")
                    st.error("Шаблон не сохранен")
                elif len(parameter_descriptions)!= len(parameters):#проверка длинны списка ключей
                    st.warning("Не указано описание для каждого параметра")
                    st.error("Шаблон не сохранен")
                else:
                    template_new_path = f"./Data/Templates/{template_name}.docx"
                    doc.save(template_new_path)#doc - название переменной, соответствует открытому файлу, загруженному пользователем
                    templates[template_name] = {
                        "description": template_description,
                        'template': template_new_path,
                    'parameters': parameter_descriptions
                    }
                    with open("./Data/формы_шаблонов.json", "w") as templates_file:
                        json.dump(templates, templates_file)
                    st.success("Шаблон сохранен")

        elif действие == "Удалить шаблоны":
            шаблоны_на_удаление = st.selectbox ("Шаблоны на удаление", [""] + list(templates.keys()))
            if шаблоны_на_удаление !="":
                удаление_шаблона = st.button("Удалить")
                if удаление_шаблона == True:
                    del templates[шаблоны_на_удаление]
                    with open("./Data/формы_шаблонов.json", "w") as templates_file:
                        json.dump(templates, templates_file)
                    st.success("Шаблон удален")

        elif действие == "Добавить новость":
            title = st.text_input(f"Наименование новости:")
            content = st.text_input(f"Новость")
            date = st.date_input("Дата новости")
             # Кнопка сохранения шаблона
            сохранение_новости = st.button("Сохранить новость")
            if сохранение_новости == True:
                if not title:
                    st.warning("Название новости не заполнено")
                    st.error("Новость не сохранена")
                elif not content:
                    st.warning("Новость не заполнена")
                    st.error("Новость не сохранена")
                elif not date:
                    st.warning("Дата не заполнена")
                    st.error("Новость не сохранена")
                else:
                    news[title] = {
                        "content": content,
                        'date': str(date)
                    }
                    with open("./Data/news.json", "w") as news_file:
                        json.dump(news, news_file)
                    st.success("Новость сохранена")

        elif действие == "Удалить новость":
            новости_на_удаление = st.selectbox ("Новости на удаление", [""] + list(news.keys()))
            if новости_на_удаление !="":
                удаление_новости = st.button("Удалить")
                if удаление_новости == True:
                    del news[новости_на_удаление]
                    with open("./Data/news.json", "w") as news_file:
                        json.dump(news, news_file)
                    st.success("Новость удалена")
# Содержимое страницы с разделом Обучение по ПОД/ФТ

# Отображение функционала только если пользователь нажал кнопку
    elif presentation_page or st.session_state.get("page") == "presentation_page":
        st.session_state.page = "presentation_page"

        st.title("Обучение по ПОД/ФТ")
        st.write("Здесь вы можете создать презентацию для проведения обучения сотрудников")

        # Функция генерации презентации
        def create_presentation(slides_data):
            prs = Presentation()
            for slide_data in slides_data:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title = slide.shapes.title
                subtitle = slide.placeholders[1]
                title.text = slide_data['title']
                subtitle.text = slide_data['subtitle']

                for i, text in enumerate(slide_data.get('text', [])):
                    left = Inches(1)
                    top = Inches(2 + i)
                    width = Inches(8.5)
                    height = Inches(1)
                    textbox = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = textbox.text_frame
                    p = text_frame.add_paragraph()
                    p.text = text
                    p.font.size = Pt(14)

                if 'image' in slide_data and slide_data['image'] is not None:
                    image_path = slide_data['image']
                    image_bytes = BytesIO(image_path.read())
                    slide.shapes.add_picture(image_bytes, Inches(1), Inches(3.5), width=Inches(4), height=Inches(3))

            pptx_file = BytesIO()
            prs.save(pptx_file)
            pptx_file.seek(0)
            return pptx_file

        # Интерфейс генерации слайдов
        st.subheader("Создайте свою презентацию")
        presentation_title = st.text_input("Название презентации", "Моя Презентация")
        presentation_subtitle = st.text_input("Подзаголовок", "Описание презентации")
        slides_data = []
        slide_count = st.number_input("Количество слайдов", min_value=1, max_value=10, step=1)

        for slide_num in range(slide_count):
            st.subheader(f"Слайд #{slide_num + 1}")
            slide_title = st.text_input(f"Заголовок для слайда #{slide_num + 1}", key=f"slide_title_{slide_num}")
            slide_subtitle = st.text_input(f"Подзаголовок для слайда #{slide_num + 1}", key=f"slide_subtitle_{slide_num}")
            slide_text = []
            num_texts = st.number_input(f"Количество текстовых блоков для слайда #{slide_num + 1}",
                                        min_value=0, max_value=5, step=1)
            for i in range(num_texts):
                text_block = st.text_input(f"Текст {i + 1} для слайда #{slide_num + 1}",
                                           key=f"slide_text_{slide_num}_{i}")
                if text_block:
                    slide_text.append(text_block)

            slide_image = st.file_uploader(f"Загрузите изображение для слайда #{slide_num + 1}",
                                           type=["jpg", "png"], key=f"slide_image_{slide_num}")

            slides_data.append({
                'title': slide_title,
                'subtitle': slide_subtitle,
                'text': slide_text,
                'image': slide_image if slide_image else None
            })

        if st.button("Сгенерировать презентацию"):
            if not slides_data:
                st.warning("Не добавлены слайды для презентации.")
            else:
                pptx_file = create_presentation(slides_data)
                st.download_button(
                    label="Скачать презентацию",
                    data=pptx_file,
                    file_name=f"{presentation_title}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )








st.markdown("<div style='height:200px;'></div>", unsafe_allow_html=True)
st.markdown("---")
st.image(logo, width=100)
st.caption("© 2025 ПВК-Конструктор")






